import re
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader

from app.services.llm_service import parse_summary_content, stream_summarize_document, summarize_document

if TYPE_CHECKING:
    from sqlalchemy.orm import Session

    from app.models.doc import Doc

def _clean_text(text: str) -> str:
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{2,}", "\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()

def _load_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def _load_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        return "\n".join([d.page_content for d in docs])
    if ext == ".docx":
        loader = Docx2txtLoader(file_path)
        docs = loader.load()
        return "\n".join([d.page_content for d in docs])
    if ext == ".txt":
        loader = TextLoader(file_path, encoding="utf-8")
        docs = loader.load()
        return "\n".join([d.page_content for d in docs])
    if ext == ".pptx":
        return _load_pptx_text(file_path)
    raise ValueError("Unsupported file type")

def parse_and_index(
    file_path: str,
    doc_id: str,
    owner_id: str,
    db: "Session | None" = None,
    doc: "Doc | None" = None,
) -> None:
    raw_text = _load_text(file_path)
    cleaned = _clean_text(raw_text)

    # 1. 大模型总结，返回结构化数据
    parsed = summarize_document(cleaned)
    if db is not None and doc is not None:
        doc.parsed_school = parsed.get("school") or ""
        doc.parsed_major = parsed.get("major") or ""
        doc.parsed_course = parsed.get("course") or ""
        doc.parsed_summary = parsed.get("summary") or ""
        kps = parsed.get("knowledgePoints")
        doc.parsed_knowledge_points = kps if isinstance(kps, list) else []
        db.commit()


def parse_and_index_stream(
    file_path: str,
    doc_id: str,
    owner_id: str,
    db: "Session | None" = None,
    doc: "Doc | None" = None,
):
    """
    流式解析文档，yield 内容片段（可用于 SSE 打字机效果）。
    结束后会将解析结果写入 doc。
    """
    raw_text = _load_text(file_path)
    cleaned = _clean_text(raw_text)

    buffer: list[str] = []
    for chunk in stream_summarize_document(cleaned):
        buffer.append(chunk)
        yield chunk

    full = "".join(buffer)
    parsed = parse_summary_content(full)
    if db is not None and doc is not None:
        doc.parsed_school = parsed.get("school") or ""
        doc.parsed_major = parsed.get("major") or ""
        doc.parsed_course = parsed.get("course") or ""
        doc.parsed_summary = parsed.get("summary") or ""
        kps = parsed.get("knowledgePoints")
        doc.parsed_knowledge_points = kps if isinstance(kps, list) else []
        db.commit()
